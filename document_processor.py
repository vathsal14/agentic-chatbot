import os
from typing import List, Dict, Any, Optional
import PyPDF2
from pptx import Presentation
from docx import Document
import pandas as pd
import numpy as np
from typing import List, Dict, Any

class DocumentProcessor:
    """Process various document formats and extract text."""
    
    @staticmethod
    def read_pdf(file_path: str) -> str:
        """Extract text from a PDF file."""
        text = []
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text.append(page.extract_text() or '')
        return '\n'.join(text)
    
    @staticmethod
    def read_pptx(file_path: str) -> str:
        """Extract text from a PowerPoint file."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    
    @staticmethod
    def read_docx(file_path: str) -> str:
        """Extract text from a Word document."""
        doc = Document(file_path)
        return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    
    @staticmethod
    def read_csv(file_path: str) -> str:
        """Extract text from a CSV file."""
        df = pd.read_csv(file_path)
        return df.to_string()
    
    @staticmethod
    def read_txt(file_path: str) -> str:
        """Read text from a plain text file."""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    def process_document(self, file_path: str) -> str:
        """Process a document based on its file extension."""
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext == '.pdf':
            return self.read_pdf(file_path)
        elif ext == '.pptx':
            return self.read_pptx(file_path)
        elif ext == '.docx':
            return self.read_docx(file_path)
        elif ext == '.csv':
            return self.read_csv(file_path)
        elif ext in ['.txt', '.md']:
            return self.read_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

class TextSplitter:
    """Split text into chunks for processing."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def split_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks."""
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = min(start + self.chunk_size, text_length)
            chunks.append(text[start:end])
            
            if end == text_length:
                break
                
            # Move back by overlap amount, but not past the start of the previous chunk
            start = max(start + self.chunk_size - self.chunk_overlap, start + 1)
        
        return chunks
