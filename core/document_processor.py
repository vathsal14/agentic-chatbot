import os
from typing import List, Dict, Any, Optional
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
import markdown
from bs4 import BeautifulSoup
import io
from typing import List
class DocumentProcessor:
    """Processes different document formats and extracts text content"""
    @staticmethod
    def process_file(file_path: str) -> Dict[str, Any]:
        """
        Process a file and return its content based on file extension
        Args:
            file_path: Path to the file to process
        Returns:
            Dict containing the extracted content and metadata
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        try:
            if file_extension == '.pdf':
                return DocumentProcessor._process_pdf(file_path)
            elif file_extension == '.docx':
                return DocumentProcessor._process_docx(file_path)
            elif file_extension == '.pptx':
                return DocumentProcessor._process_pptx(file_path)
            elif file_extension == '.csv':
                return DocumentProcessor._process_csv(file_path)
            elif file_extension in ['.txt', '.md']:
                return DocumentProcessor._process_text(file_path)
            else:
                raise ValueError(f"Unsupported file format: {file_extension}")
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "file_path": file_path
            }
    @staticmethod
    def process_document(file_path: str, metadata: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        Process a document and return its content with metadata
        Args:
            file_path: Path to the file to process
            metadata: Additional metadata to include with the document
        Returns:
            Dict containing the extracted content, metadata, and file information
        """
        if metadata is None:
            metadata = {}
        result = DocumentProcessor.process_file(file_path)
        file_info = {
            "file_path": file_path,
            "file_name": os.path.basename(file_path),
            "file_extension": os.path.splitext(file_path)[1].lower(),
            "file_size": os.path.getsize(file_path),
            "last_modified": os.path.getmtime(file_path)
        }
        result_metadata = {
            **file_info,
            **metadata
        }
        if "metadata" in result:
            result["metadata"].update(result_metadata)
        else:
            result["metadata"] = result_metadata
        return result
    @staticmethod
    def chunk_text(
        content: str,
        chunk_size: int = 1000,
        overlap: int = 200,
        separator: str = "\n"
    ) -> List[Dict[str, Any]]:
        """
        Split text into chunks of specified size with overlap
        Args:
            content: The text content to chunk
            chunk_size: Maximum size of each chunk in characters
            overlap: Number of characters to overlap between chunks
            separator: Separator to use when splitting text
        Returns:
            List of chunks with text and position information
        """
        if not content:
            return []
        paragraphs = content.split(separator)
        chunks = []
        current_chunk = []
        current_length = 0
        start_pos = 0
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            para_length = len(para)
            if current_length + para_length > chunk_size and current_chunk:
                chunk_text = separator.join(current_chunk)
                end_pos = start_pos + len(chunk_text)
                chunks.append({
                    "text": chunk_text,
                    "start": start_pos,
                    "end": end_pos
                })
                overlap_start = max(0, len(current_chunk) - overlap // (len(para) // max(1, len(current_chunk))))
                current_chunk = current_chunk[overlap_start:]
                current_length = len(separator.join(current_chunk))
                start_pos = end_pos - current_length
            current_chunk.append(para)
            current_length += para_length + len(separator)
        if current_chunk:
            chunk_text = separator.join(current_chunk)
            end_pos = start_pos + len(chunk_text)
            chunks.append({
                "text": chunk_text,
                "start": start_pos,
                "end": end_pos
            })
        return chunks
class TextSplitter:
    """Split text into chunks for processing."""
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    def split_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks."""
        if not text:
            return []
        chunks = []
        start = 0
        text_length = len(text)
        while start < text_length:
            end = min(start + self.chunk_size, text_length)
            chunks.append(text[start:end])
            if end == text_length:
                break
            start = max(start + self.chunk_size - self.chunk_overlap, start + 1)
        return chunks
    @staticmethod
    def _process_pdf(file_path: str) -> Dict[str, Any]:
        """Extract text from PDF file"""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n\n"
        return {
            "success": True,
            "content": text.strip(),
            "file_type": "pdf",
            "page_count": len(pdf_reader.pages)
        }
    @staticmethod
    def _process_docx(file_path: str) -> Dict[str, Any]:
        """Extract text from DOCX file"""
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return {
            "success": True,
            "content": "\n\n".join(full_text).strip(),
            "file_type": "docx"
        }
    @staticmethod
    def _process_pptx(file_path: str) -> Dict[str, Any]:
        """Extract text from PPTX file"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return {
            "success": True,
            "content": "\n\n".join(text).strip(),
            "file_type": "pptx",
            "slide_count": len(prs.slides)
        }
    @staticmethod
    def _process_csv(file_path: str) -> Dict[str, Any]:
        """Extract content from CSV file"""
        df = pd.read_csv(file_path)
        csv_content = df.to_string(index=False)
        return {
            "success": True,
            "content": csv_content,
            "file_type": "csv",
            "rows": len(df),
            "columns": list(df.columns)
        }
    @staticmethod
    def _process_text(file_path: str) -> Dict[str, Any]:
        """Extract text from TXT or MD file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        if file_path.lower().endswith('.md'):
            html = markdown.markdown(content)
            soup = BeautifulSoup(html, 'html.parser')
            content = soup.get_text()
        return {
            "success": True,
            "content": content.strip(),
            "file_type": "markdown" if file_path.lower().endswith('.md') else "text"
        }
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
        """
        Split text into overlapping chunks
        Args:
            text: The text to chunk
            chunk_size: Maximum size of each chunk
            overlap: Number of characters to overlap between chunks
        Returns:
            List of chunks with metadata
        """
        if not text:
            return []
        chunks = []
        start = 0
        text_length = len(text)
        while start < text_length:
            end = min(start + chunk_size, text_length)
            chunk = text[start:end]
            chunks.append({
                "text": chunk,
                "start": start,
                "end": end,
                "length": len(chunk)
            })
            if end == text_length:
                break
            start = end - overlap
        return chunks